"""
Generate israel_researcher_overview.pptx
Dark-themed, illustrated flow presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Palette ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT1  = RGBColor(0x00, 0xC8, 0xFF)   # cyan
ACCENT2  = RGBColor(0x00, 0xE5, 0x96)   # teal-green
ACCENT3  = RGBColor(0xFF, 0xB3, 0x00)   # amber
ACCENT4  = RGBColor(0xFF, 0x4D, 0x6D)   # coral-red
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xB0, 0xC4, 0xDE)
DGRAY    = RGBColor(0x1E, 0x30, 0x45)
MGRAY    = RGBColor(0x2A, 0x40, 0x58)

W = 13.33   # slide width  (inches, widescreen)
H = 7.5     # slide height (inches)

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=BG):
    sp = slide.shapes.add_shape(1, 0, 0, Inches(W), Inches(H))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp

def rect(slide, x, y, w, h, color, radius=False):
    shape_type = 5 if radius else 1   # 5=rounded, 1=rect
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp

def label(slide, text, x, y, w, h,
          size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def box_label(slide, text, x, y, w, h, bg_color, fg_color=WHITE,
              size=13, bold=True, align=PP_ALIGN.CENTER, radius=False):
    rect(slide, x, y, w, h, bg_color, radius=radius)
    label(slide, text, x, y, w, h, size=size, bold=bold, color=fg_color,
          align=align)

def hline(slide, x, y, w, color=ACCENT1, thick=2):
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(thick))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()

def arrow_right(slide, x, y, length=0.5, color=LGRAY):
    """Horizontal arrow using a rectangle + triangle-ish label."""
    rect(slide, x, y+0.045, length-0.15, 0.06, color)
    # arrowhead triangle (chevron via text)
    label(slide, "▶", x+length-0.2, y-0.02, 0.25, 0.18,
          size=11, color=color, align=PP_ALIGN.LEFT)

def arrow_down(slide, x, y, length=0.35, color=LGRAY):
    rect(slide, x+0.04, y, 0.06, length-0.1, color)
    label(slide, "▼", x-0.01, y+length-0.18, 0.2, 0.18,
          size=11, color=color, align=PP_ALIGN.LEFT)

def bullet_list(slide, items, x, y, w, spacing=0.32,
                size=13, icon="▸", icon_color=ACCENT2, text_color=LGRAY):
    for i, item in enumerate(items):
        label(slide, icon, x, y + i*spacing, 0.25, 0.3,
              size=size, color=icon_color, bold=True)
        label(slide, item, x+0.28, y + i*spacing, w-0.28, 0.3,
              size=size, color=text_color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)

# gradient-like accent strip
rect(s, 0, 0, W, 0.08, ACCENT1)
rect(s, 0, H-0.08, W, 0.08, ACCENT2)

# big title
label(s, "Israel Researcher", 1.2, 1.5, 10.8, 1.5,
      size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
label(s, "AI-Powered TASE Market Intelligence System",
      1.2, 3.0, 10.8, 0.7, size=22, color=ACCENT1, align=PP_ALIGN.CENTER)

hline(s, 2.5, 3.9, 8.3, ACCENT2, thick=3)

label(s, "Automated • Multi-Agent • Real-Time • Telegram Alerts",
      1.2, 4.15, 10.8, 0.5, size=15, color=LGRAY, align=PP_ALIGN.CENTER,
      italic=True)

# stats boxes
for i, (val, lbl) in enumerate([
    ("15 min", "Cycle"), ("7 Agents", "Parallel"), ("40+", "Signal Types"),
    ("130+", "Convergence Rules"), ("3 Phases", "Pipeline"),
]):
    bx = 1.0 + i * 2.3
    rect(s, bx, 5.1, 2.0, 1.5, MGRAY, radius=True)
    label(s, val,  bx, 5.25, 2.0, 0.6, size=26, bold=True, color=ACCENT1,
          align=PP_ALIGN.CENTER)
    label(s, lbl,  bx, 5.85, 2.0, 0.4, size=13, color=LGRAY,
          align=PP_ALIGN.CENTER)

label(s, "Brokai © 2026", 0, H-0.4, W, 0.3,
      size=10, color=MGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SYSTEM OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, ACCENT1)

label(s, "System Overview", 0.3, 0.15, 12, 0.65,
      size=32, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, ACCENT2, thick=2)

# three columns
cols = [
    ("📥  DATA SOURCES", ACCENT1, [
        "Maya TASE regulatory filings",
        "Israeli RSS news (Ynet, Walla, Maariv)",
        "Chrome-scraped: Globes, Calcalist, TheMarker",
        "Google News RSS (per-ticker)",
        "yfinance — prices, volume, technicals",
        "Dual-listed US overnight moves",
        "Global macro (S&P500, VIX, USD/ILS)",
    ]),
    ("⚙️  PROCESSING", ACCENT3, [
        "8 technical anomaly detectors",
        "Signal convergence scoring engine",
        "40+ signal types, 130+ multiplier pairs",
        "Earnings proximity gradient",
        "LLM web news extraction (GPT-4o)",
        "Deep financial analysis (RSI, MA, rev.)",
        "Per-stock memory & history cache",
    ]),
    ("📤  OUTPUTS", ACCENT2, [
        "Quick alert every 15 min (Telegram)",
        "Daily summary at 17:00",
        "Weekly Stock of the Week (Thursday)",
        "Excel memory backup (3-sheet workbook)",
        "Dedup: no repeat picks within same week",
        "Score 0-100 with calibrated rubric",
        "Tier labels: BUY / WATCH / MONITOR",
    ]),
]

for ci, (title, col, items) in enumerate(cols):
    cx = 0.25 + ci * 4.35
    rect(s, cx, 1.05, 4.1, 5.9, MGRAY, radius=True)
    label(s, title, cx+0.1, 1.15, 3.9, 0.5,
          size=13, bold=True, color=col, align=PP_ALIGN.LEFT)
    hline(s, cx+0.1, 1.65, 3.9, col, thick=1)
    bullet_list(s, items, cx+0.1, 1.75, 3.9,
                spacing=0.7, size=11.5, icon_color=col)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — FULL PIPELINE FLOW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, ACCENT1)

label(s, "Full Pipeline — Every 15 Minutes", 0.3, 0.15, 12, 0.65,
      size=32, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, ACCENT2, thick=2)

# Flow boxes top-to-bottom, left-to-right
phases = [
    ("PHASE 1\nData Collection", BG,   ACCENT1, 0.3,  1.05, 12.7, 1.1),
    ("PHASE 2\nSector Agents",   BG,   ACCENT3, 0.3,  2.8,  12.7, 2.6),
    ("PHASE 3\nArbitration",     BG,   ACCENT2, 0.3,  6.05, 12.7, 0.75),
]

# Phase 1 box
rect(s, 0.3, 1.05, 12.7, 1.1, DGRAY, radius=True)
rect(s, 0.3, 1.05, 0.18, 1.1, ACCENT1)
label(s, "PHASE 1  —  Sequential Data Collection", 0.6, 1.1, 8, 0.4,
      size=14, bold=True, color=ACCENT1)
p1_items = ["Maya Filings", "Earnings Calendar", "Dual-listed US", "Israeli RSS", "Chrome News", "Macro Snapshot", "Sector Rotation"]
for i, item in enumerate(p1_items):
    bx = 0.55 + i * 1.75
    rect(s, bx, 1.55, 1.6, 0.45, MGRAY, radius=True)
    label(s, item, bx, 1.58, 1.6, 0.42, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Phase 2 box
rect(s, 0.3, 2.8, 12.7, 2.6, DGRAY, radius=True)
rect(s, 0.3, 2.8, 0.18, 2.6, ACCENT3)
label(s, "PHASE 2  —  Parallel Sector Agents  (4 workers × 7 agents)", 0.6, 2.85, 10, 0.4,
      size=14, bold=True, color=ACCENT3)

agents = [
    ("🏦 Banks", ACCENT1),
    ("🛡 Tech/Defense", ACCENT2),
    ("⚡ Energy", ACCENT3),
    ("💊 Pharma", ACCENT4),
    ("🏢 Real Estate", RGBColor(0xA0,0x6F,0xFF)),
    ("📡 Telecom/Consumer", RGBColor(0xFF,0x6F,0xA0)),
    ("🔍 Discovery", LGRAY),
]
for i, (name, col) in enumerate(agents):
    ax = 0.55 + i * 1.75
    rect(s, ax, 3.3, 1.65, 0.55, MGRAY, radius=True)
    label(s, name, ax, 3.32, 1.65, 0.55, size=10.5, bold=True,
          color=col, align=PP_ALIGN.CENTER)

# sub-steps inside each agent
steps = ["Filter signals", "8 Tech detectors", "Sector macro", "Convergence", "Web news (LLM)", "Deep analyze", "Sector LLM"]
for i, step in enumerate(steps):
    ax = 0.55 + i * 1.75
    label(s, step, ax, 3.9, 1.65, 0.4, size=9, color=LGRAY, align=PP_ALIGN.CENTER)

# convergence arrow down to Phase 2 output
rect(s, 4.5, 4.35, 4.3, 0.45, RGBColor(0x15,0x25,0x38), radius=True)
label(s, "↓  all sector portfolios  ↓", 4.5, 4.38, 4.3, 0.42,
      size=12, color=LGRAY, align=PP_ALIGN.CENTER)

# Phase 3 box
rect(s, 0.3, 6.05, 12.7, 0.75, DGRAY, radius=True)
rect(s, 0.3, 6.05, 0.18, 0.75, ACCENT2)
label(s, "PHASE 3  —  Manager / CIO LLM Arbitration", 0.6, 6.1, 7, 0.35,
      size=14, bold=True, color=ACCENT2)
label(s, "Stock of the Week  +  Runners-Up  →  Telegram Alert  →  Excel Log",
      0.6, 6.47, 12, 0.3, size=11, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PHASE 1 DEEP DIVE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, ACCENT1)

label(s, "Phase 1 — Cross-Sector Data Collection", 0.3, 0.15, 12, 0.65,
      size=28, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, ACCENT1, thick=2)
label(s, "Runs once per cycle — shared results passed to all 7 sector agents",
      0.3, 0.9, 12.7, 0.35, size=13, color=LGRAY, italic=True)

sources = [
    ("Maya Filings", ACCENT1, "Playwright browser bypasses Incapsula WAF\n→ 100 latest regulatory filings\n→ signal types: contract, IPO, earnings,\n   institutional, buyback, dividend, management"),
    ("Earnings Calendar", ACCENT3, "POST /api/v1/corporate-actions/...\n→ events within 10 days\n→ earnings_calendar signal\n→ triggers gradient bonus (+80 pts at dte=0)"),
    ("Dual-Listed US Moves", ACCENT2, "10 stocks: TEVA, NICE, ICL, ESLT...\n→ yfinance overnight price change\n→ ≥2% move → dual_listed_move signal\n→ Most reliable TASE-open predictor"),
    ("Israeli RSS News", RGBColor(0xA0,0x6F,0xFF), "Ynet, Walla, Maariv RSS feeds\n→ Hebrew company name regex matching\n→ SignalEnricher upgrades signal type\n→ word-boundary keyword detection"),
    ("Chrome News (NEW)", ACCENT4, "Playwright scrapes Globes, Calcalist,\nTheMarker — reuses Maya browser session\n→ JS-rendered pages, no bot blocking\n→ Same items_to_signals() pipeline"),
    ("Macro Snapshot", LGRAY, "TA-125, S&P500, VIX, USD/ILS, Nasdaq\n→ formatted text string\n→ injected into every LLM call\n→ adjusts scores ±5 to ±10"),
]

for i, (title, col, detail) in enumerate(sources):
    cx = 0.25 + (i % 3) * 4.35
    cy = 1.4 + (i // 3) * 2.8
    rect(s, cx, cy, 4.1, 2.55, MGRAY, radius=True)
    rect(s, cx, cy, 4.1, 0.4, col, radius=True)
    rect(s, cx, cy+0.4, 4.1, 0.03, col)   # fill corner
    label(s, title, cx+0.1, cy+0.03, 3.9, 0.38,
          size=13, bold=True, color=BG, align=PP_ALIGN.LEFT)
    label(s, detail, cx+0.12, cy+0.48, 3.85, 1.95,
          size=10.5, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SIGNAL SCORING & CONVERGENCE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, ACCENT3)

label(s, "Signal Scoring & Convergence Engine", 0.3, 0.15, 12, 0.65,
      size=28, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, ACCENT3, thick=2)

# Formula box
rect(s, 0.3, 1.0, 12.7, 1.0, MGRAY, radius=True)
label(s, "final_score  =  min(100,  Σ base_score  ×  best_multiplier  ×  3-category-boost  +  earnings_gradient)",
      0.5, 1.1, 12.3, 0.6, size=14, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

# Left: top signal scores
rect(s, 0.3, 2.2, 3.8, 4.8, DGRAY, radius=True)
label(s, "TOP BASE SCORES", 0.4, 2.28, 3.6, 0.4,
      size=12, bold=True, color=ACCENT3)
scores = [
    ("maya_ipo", 50), ("maya_contract", 45), ("new_contract", 45),
    ("regulatory_approval", 42), ("maya_institutional", 40),
    ("breakout", 35), ("dual_listed_move", 35), ("earnings (web)", 35),
    ("volume_spike", 30), ("price_move", 25),
]
for i, (sig, sc) in enumerate(scores):
    cy = 2.75 + i * 0.4
    rect(s, 0.35, cy, 2.4, 0.33, MGRAY, radius=True)
    label(s, sig, 0.45, cy+0.03, 2.3, 0.3, size=10, color=WHITE)
    rect(s, 2.8, cy, 0.9, 0.33, ACCENT3, radius=True)
    label(s, str(sc), 2.8, cy+0.03, 0.9, 0.3, size=10, bold=True,
          color=BG, align=PP_ALIGN.CENTER)

# Middle: top multipliers
rect(s, 4.3, 2.2, 4.5, 4.8, DGRAY, radius=True)
label(s, "TOP MULTIPLIERS", 4.4, 2.28, 4.3, 0.4,
      size=12, bold=True, color=ACCENT1)
mults = [
    ("earnings_cal + volume_spike", "2.5×"),
    ("low_reversal + earnings_cal", "2.4×"),
    ("oversold_bounce + earnings_cal", "2.3×"),
    ("earnings + volume_spike", "2.3×"),
    ("breakout + earnings_cal", "2.3×"),
    ("rel_strength + earnings_cal", "2.2×"),
    ("earnings_cal + new_contract", "2.2×"),
    ("dual_listed + earnings_cal", "2.2×"),
    ("institutional + volume_spike", "1.8×"),
    ("3+ categories → bonus", "×1.3"),
]
for i, (pair, mult) in enumerate(mults):
    cy = 2.75 + i * 0.4
    rect(s, 4.35, cy, 3.35, 0.33, MGRAY, radius=True)
    label(s, pair, 4.45, cy+0.03, 3.25, 0.3, size=9.5, color=WHITE)
    rect(s, 7.75, cy, 0.65, 0.33, ACCENT1, radius=True)
    label(s, mult, 7.75, cy+0.03, 0.65, 0.3, size=10, bold=True,
          color=BG, align=PP_ALIGN.CENTER)

# Right: earnings gradient + calibration
rect(s, 9.0, 2.2, 4.0, 2.3, DGRAY, radius=True)
label(s, "EARNINGS GRADIENT BONUS", 9.1, 2.28, 3.8, 0.4,
      size=12, bold=True, color=ACCENT2)
grad = [("dte = 0 (today)", "+80"), ("dte = 1", "+70"), ("dte = 2", "+60"),
        ("dte = 3", "+45"), ("dte ≤ 7", "+25"), ("dte ≤ 14", "+12")]
for i, (d, bonus) in enumerate(grad):
    cy = 2.72 + i * 0.3
    label(s, d,     9.1,  cy, 2.5, 0.28, size=10, color=LGRAY)
    label(s, bonus, 11.6, cy, 0.9, 0.28, size=10, bold=True, color=ACCENT2,
          align=PP_ALIGN.RIGHT)

rect(s, 9.0, 4.7, 4.0, 2.3, DGRAY, radius=True)
label(s, "SCORE CALIBRATION", 9.1, 4.78, 3.8, 0.4,
      size=12, bold=True, color=ACCENT4)
calib = [
    ("88-100", "Exceptional — all 4 criteria"),
    ("72-87",  "Strong — hard catalyst + tech"),
    ("55-71",  "Moderate — 1 element missing"),
    ("40-54",  "Watch — single signal"),
    ("<40",    "Monitor — weak signal"),
]
for i, (rng, desc) in enumerate(calib):
    cy = 5.18 + i * 0.36
    rect(s, 9.05, cy, 1.0, 0.3, ACCENT4, radius=True)
    label(s, rng, 9.05, cy+0.02, 1.0, 0.28, size=9, bold=True,
          color=BG, align=PP_ALIGN.CENTER)
    label(s, desc, 10.1, cy+0.02, 2.75, 0.28, size=9.5, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SECTOR AGENTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, ACCENT3)

label(s, "Phase 2 — Sector Agents (Parallel)", 0.3, 0.15, 12, 0.65,
      size=28, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, ACCENT3, thick=2)

# per-agent flow (left column)
rect(s, 0.25, 1.0, 4.2, 6.3, DGRAY, radius=True)
label(s, "PER-AGENT FLOW", 0.4, 1.05, 4.0, 0.4,
      size=13, bold=True, color=ACCENT3)
steps = [
    ("1", "Filter pre-fetched signals for sector tickers", ACCENT1),
    ("2", "MarketAnomalyDetector — 8 technical detectors", ACCENT1),
    ("3", "Sector macro signals (oil / VIX / rates / peers)", ACCENT3),
    ("4", "Preliminary ConvergenceEngine grouping", ACCENT3),
    ("5", "Filter to tickers with final_score > 0 (discovery-first)", ACCENT2),
    ("6", "Google News + LLM extraction → web signals (top-5)", ACCENT2),
    ("7", "Re-run convergence with enriched signals", ACCENT2),
    ("8", "DeepStockAnalyzer: RSI, MA, revenue growth (top-8)", ACCENT4),
    ("9", "Sector LLM → full ranked portfolio", ACCENT4),
]
for i, (num, desc, col) in enumerate(steps):
    cy = 1.55 + i * 0.62
    rect(s, 0.35, cy, 0.38, 0.38, col, radius=True)
    label(s, num, 0.35, cy+0.02, 0.38, 0.36, size=12, bold=True,
          color=BG, align=PP_ALIGN.CENTER)
    label(s, desc, 0.82, cy+0.04, 3.5, 0.35, size=10.5, color=LGRAY)

# 8 technical detectors
rect(s, 4.65, 1.0, 4.0, 4.1, DGRAY, radius=True)
label(s, "8 TECHNICAL DETECTORS", 4.75, 1.05, 3.8, 0.4,
      size=13, bold=True, color=ACCENT1)
detectors = [
    ("volume_spike",         "Volume >2.5× 20-day avg"),
    ("price_move",           "Daily move >3.5%"),
    ("breakout",             "Within 3% of 52-week high + vol"),
    ("ma_crossover",         "MA-20 crosses above MA-50"),
    ("oversold_bounce",      "RSI <32 + rising volume"),
    ("low_reversal",         "Within 5% of 52-week low + bounce"),
    ("consecutive_momentum", "4+ consecutive up days"),
    ("relative_strength",    "Outperforming TA-125 by >5%"),
]
for i, (sig, desc) in enumerate(detectors):
    cy = 1.55 + i * 0.45
    rect(s, 4.7, cy, 1.9, 0.38, MGRAY, radius=True)
    label(s, sig, 4.75, cy+0.03, 1.85, 0.34, size=9, bold=True, color=ACCENT1)
    label(s, desc, 6.65, cy+0.03, 1.85, 0.34, size=9, color=LGRAY)

# sector-specific macro signals
rect(s, 4.65, 5.3, 4.0, 2.0, DGRAY, radius=True)
label(s, "SECTOR MACRO SIGNALS", 4.75, 5.35, 3.8, 0.4,
      size=13, bold=True, color=ACCENT3)
macros = [
    ("Banks",       "KBE/XLF ETF move, IL10Y yield"),
    ("TechDefense", "LMT/RTX peers, shekel, VIX defense"),
    ("Energy",      "WTI, Brent, Natural Gas prices"),
    ("Pharma",      "XBI/IBB biotech ETF move"),
    ("Real Estate", "VNQ/IYR REIT ETFs, shekel"),
    ("Telecom",     "XLP/IYZ consumer ETFs"),
]
for i, (sec, sig) in enumerate(macros):
    cy = 5.82 + i * 0.24
    label(s, f"• {sec}:", 4.75, cy, 1.2, 0.22, size=9, bold=True, color=ACCENT3)
    label(s, sig, 5.9, cy, 2.65, 0.22, size=9, color=LGRAY)

# Discovery agent box
rect(s, 8.85, 1.0, 4.2, 6.3, DGRAY, radius=True)
label(s, "🔍 DISCOVERY AGENT", 8.95, 1.05, 4.0, 0.4,
      size=13, bold=True, color=LGRAY)
bullet_list(s, [
    "Reads full Maya company cache (~500 companies)",
    "Validates each {SYMBOL}.TA on Yahoo Finance",
    "Caches results: valid=30d, invalid=7d TTL",
    "Max 25 new validations per cycle",
    "Priority: tickers from today's Maya filings",
    "Covers IPOs, TA-SME, uncovered tickers",
    "Runs same 8-detector technical scan",
    "No sector macro signals (too diverse)",
    "Fills the gap between listing and coverage",
], x=9.0, y=1.55, w=3.9,
   spacing=0.58, size=10.5,
   icon="→", icon_color=LGRAY, text_color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PHASE 3: MANAGER LLM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, ACCENT2)

label(s, "Phase 3 — Manager / CIO Arbitration", 0.3, 0.15, 12, 0.65,
      size=28, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, ACCENT2, thick=2)

# Input box
rect(s, 0.3, 1.05, 3.5, 5.8, DGRAY, radius=True)
label(s, "INPUT", 0.45, 1.1, 3.2, 0.38, size=13, bold=True, color=ACCENT1)
bullet_list(s, [
    "Full sector portfolio from each of 7 agents",
    "Each stock: tier, score, rationale, catalyst, keywords",
    "Macro snapshot (TA-125, VIX, USD/ILS, SP500)",
    "Sector rotation context (BULL+/BEAR- labels)",
    "Stock memory: historical signals & analyst notes",
], x=0.45, y=1.56, w=3.25,
   spacing=0.62, size=11, icon_color=ACCENT1)

# CIO rules
rect(s, 4.05, 1.05, 5.4, 5.8, DGRAY, radius=True)
label(s, "CIO DECISION RULES", 4.2, 1.1, 5.1, 0.38,
      size=13, bold=True, color=ACCENT3)
rules = [
    "No 2 picks from same sector (unless score >90)",
    "Prioritize sector with strongest macro tailwind",
    "Balance large-cap with at least 1 mid/small-cap",
    "Best pick must have: catalyst + technicals + macro",
    "Sector conviction (tier=BUY, BULL+) > raw score",
    "Diversify: not all cyclical or all defensive",
    "REJECT stocks with sector-only macro signals",
    "Score calibrated: most winners score 70–82",
    "Score 88-100 = exceptional, at most 1×/month",
]
for i, rule in enumerate(rules):
    cy = 1.55 + i * 0.57
    rect(s, 4.12, cy, 0.35, 0.35, ACCENT3, radius=True)
    label(s, str(i+1), 4.12, cy+0.02, 0.35, 0.33, size=10, bold=True,
          color=BG, align=PP_ALIGN.CENTER)
    label(s, rule, 4.53, cy+0.04, 4.75, 0.33, size=10.5, color=LGRAY)

# Output box
rect(s, 9.65, 1.05, 3.4, 5.8, DGRAY, radius=True)
label(s, "OUTPUT JSON", 9.8, 1.1, 3.1, 0.38, size=13, bold=True, color=ACCENT2)
out_items = [
    'stock_of_the_week:\n  ticker, name, score, tier,\n  full_rationale (4-6 sentences),\n  key_catalyst, technical_setup,\n  main_risk, sector, keywords',
    'runners_up:\n  [{ticker, name, score, tier,\n    summary, key_catalyst, sector}]',
    'macro_context:\n  2-3 sentence macro summary',
    'week_theme:\n  dominant theme this week',
    'sector_in_focus:\n  strongest sector',
]
cy = 1.58
for item in out_items:
    lines = item.split('\n')
    label(s, "▸ " + lines[0], 9.8, cy, 3.1, 0.3, size=10.5,
          bold=True, color=ACCENT2)
    for extra in lines[1:]:
        cy += 0.25
        label(s, extra, 9.95, cy, 3.0, 0.25, size=9.5, color=LGRAY, italic=True)
    cy += 0.45

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ALERTS & DEDUPLICATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, ACCENT4)

label(s, "Alert System & Deduplication", 0.3, 0.15, 12, 0.65,
      size=28, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, ACCENT4, thick=2)

# 3 alert types
types = [
    ("⚡  Quick Alert", ACCENT1, "Every 15-minute cycle",
     ["Top 3 stocks from manager arbitration",
      "Only new tickers (not sent this week)",
      "Score + tier + technical setup + main risk",
      "Key catalyst clearly stated",
      "Themes / keywords listed"]),
    ("📊  Daily Summary", ACCENT3, "Once daily at 17:00 Israel time",
     ["Top 3 from today's arbitration",
      "Filters already-alerted tickers first",
      "Falls back to full ranked list if all filtered",
      "Header: '--- Daily TASE Summary ---'",
      "Same rich format as quick alert"]),
    ("🏆  Stock of the Week", ACCENT2, "Every Thursday at 17:00",
     ["Full 4-6 sentence research rationale",
      "Catalyst, technical setup, main risk",
      "2 runners-up with summary",
      "Week theme + sector in focus",
      "Macro context paragraph"]),
]
for ci, (title, col, timing, items) in enumerate(types):
    cx = 0.3 + ci * 4.35
    rect(s, cx, 1.1, 4.1, 5.3, MGRAY, radius=True)
    rect(s, cx, 1.1, 4.1, 0.55, col, radius=True)
    rect(s, cx, 1.55, 4.1, 0.08, col)
    label(s, title, cx+0.15, 1.13, 3.8, 0.38, size=14, bold=True, color=BG)
    label(s, timing, cx+0.15, 1.62, 3.8, 0.3, size=11, color=col, italic=True)
    bullet_list(s, items, cx+0.1, 2.0, 3.9,
                spacing=0.65, size=11, icon_color=col, text_color=LGRAY)

# dedup rules
rect(s, 0.3, 6.55, 12.7, 0.8, DGRAY, radius=True)
label(s, "Deduplication Rules:", 0.5, 6.6, 2.4, 0.35,
      size=12, bold=True, color=ACCENT4)
dedup = [
    "Same ticker not quick-alerted twice in same week (tracked in state + Excel 'Sent Alerts' sheet)",
    "Weekly pick: if same stock as last week AND no new signals this week → promote first runner-up",
    "All sent alerts logged to Excel with timestamp, week, score, catalyst — survives full state reset",
]
for i, d in enumerate(dedup):
    label(s, f"• {d}", 2.95 + i*0.0, 6.6 + i*0.22, 9.9, 0.22,
          size=9.5, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MEMORY & STATE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, RGBColor(0xA0,0x6F,0xFF))

label(s, "Memory, State & Persistence", 0.3, 0.15, 12, 0.65,
      size=28, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, RGBColor(0xA0,0x6F,0xFF), thick=2)

# JSON state
rect(s, 0.3, 1.1, 4.2, 5.9, DGRAY, radius=True)
label(s, "STATE FILE  (israel_researcher_state.json)", 0.4, 1.15, 4.0, 0.38,
      size=11, bold=True, color=RGBColor(0xA0,0x6F,0xFF))
state_items = [
    ("seen_maya_report_ids", "Maya filing dedup"),
    ("seen_signal_keys",     "Technical signal dedup (today only)"),
    ("weekly_signals",       "Accumulated signals this week (cap 500)"),
    ("tase_company_cache",   "All TASE companies (24h TTL)"),
    ("ticker_validation_cache", "Yahoo Finance validity (30d/7d TTL)"),
    ("stock_memory",         "Per-ticker fundamentals + history"),
    ("alerted_this_week",    "Tickers sent in quick alerts this week"),
    ("last_weekly_pick",     "Last week's stock of the week ticker"),
    ("last_daily_report",    "Timestamp of last daily summary"),
    ("last_weekly_report",   "Timestamp of last weekly report"),
]
for i, (key, desc) in enumerate(state_items):
    cy = 1.62 + i * 0.52
    rect(s, 0.38, cy, 2.05, 0.38, MGRAY, radius=True)
    label(s, key, 0.43, cy+0.04, 2.0, 0.32, size=8.5, color=ACCENT1)
    label(s, desc, 2.5, cy+0.04, 1.85, 0.32, size=9, color=LGRAY)

# Memory store
rect(s, 4.7, 1.1, 4.3, 5.9, DGRAY, radius=True)
label(s, "PER-STOCK MEMORY  (StockMemoryManager)", 4.8, 1.15, 4.1, 0.38,
      size=11, bold=True, color=ACCENT2)
bullet_list(s, [
    "fundamentals: RSI, MA-20/50, market cap, revenue growth  (7-day TTL)",
    "signal_history: last 10 cycles of signals + scores",
    "analyst_notes: LLM-written rationale from prior cycles",
    "recent_news: top-3 article headlines as pipe-separated string",
    "consecutive_active: how many cycles this week the stock appeared",
    "Prunes entries inactive >30 days each cycle",
    "Context injected into sector LLM each cycle",
    "Cached fundamentals skip yfinance API calls (quota saving)",
], x=4.85, y=1.62, w=4.05,
   spacing=0.68, size=10, icon_color=ACCENT2, text_color=LGRAY)

# Excel
rect(s, 9.2, 1.1, 3.9, 5.9, DGRAY, radius=True)
label(s, "EXCEL BACKUP  (israel_researcher_memory.xlsx)", 9.3, 1.15, 3.7, 0.38,
      size=11, bold=True, color=ACCENT3)
sheets = [
    ("Sheet 1: Active Memory",
     "Full mirror of stock_memory\nOverwritten each cycle\nRestore source if state.json deleted"),
    ("Sheet 2: Research Log",
     "Append-only buy/watch picks\nDedup by (Date, Ticker)\nSorted newest-first by score"),
    ("Sheet 3: Sent Alerts",
     "Every Telegram alert logged\nTimestamp, week, type, ticker,\nscore, key catalyst\nRestores dedup state after reset"),
]
cy = 1.62
for title, detail in sheets:
    rect(s, 9.28, cy, 3.72, 1.7, MGRAY, radius=True)
    label(s, title, 9.38, cy+0.05, 3.52, 0.35, size=11, bold=True, color=ACCENT3)
    label(s, detail, 9.38, cy+0.42, 3.52, 1.2, size=9.5, color=LGRAY)
    cy += 1.85

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — KNOWN LIMITATIONS & ARCHITECTURE NOTES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, W, 0.08, LGRAY)

label(s, "Architecture Notes & Known Limitations", 0.3, 0.15, 12, 0.65,
      size=28, bold=True, color=WHITE)
hline(s, 0.3, 0.85, 12.7, LGRAY, thick=2)

cols2 = [
    ("KEY DESIGN DECISIONS", ACCENT1, [
        "Playwright reused for both Maya + Chrome news — avoids launching 2 browsers",
        "Discovery-first: only tickers WITH signals surface — no noise from silent stocks",
        "Web news is the bridge between Maya pseudo-tickers (TASE92) and real .TA symbols",
        "ThreadPoolExecutor (4 workers) — safe for yfinance + OpenAI API rate limits",
        "Weekly accumulator deduplication by (ticker, signal_type, date) — prevents repeat",
        "Earnings gradient makes upcoming earnings the single biggest score multiplier",
        "Memory cache skips yfinance for fresh fundamentals — major API quota saving",
        "Sector-only signals (shekel_move, oil_correlation) capped at score 52 to prevent HARL-type false positives",
    ]),
    ("KNOWN LIMITATIONS", ACCENT4, [
        "Maya pseudo-tickers (TASE{id}) never converge with .TA technical signals — web news fills gap",
        "yfinance calls have no timeout — Yahoo slowness can block an agent for seconds",
        "Google News rate-limits (429) handled with 1 retry + 3s backoff only",
        "Hebrew company name regex matching often produces GENERAL ticker (name variants)",
        "State JSON has no backup strategy — Excel provides partial recovery only",
        "Delisted stocks must be manually removed: SPNS (Dec 2025), MGIC (Feb 2026)",
        "^TA35.TA invalid on Yahoo Finance — must use ^TA125.TA",
        "Clal Insurance: CLIS.TA not CLAL.TA",
    ]),
]

for ci, (title, col, items) in enumerate(cols2):
    cx = 0.3 + ci * 6.6
    rect(s, cx, 1.05, 6.4, 6.2, DGRAY, radius=True)
    label(s, title, cx+0.15, 1.1, 6.1, 0.38, size=13, bold=True, color=col)
    hline(s, cx+0.1, 1.5, 6.2, col, thick=1)
    bullet_list(s, items, cx+0.1, 1.6, 6.2,
                spacing=0.68, size=10.5, icon_color=col, text_color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
from pathlib import Path as _Path
_docs = _Path(__file__).parent.parent / "docs"
_docs.mkdir(exist_ok=True)
out = str(_docs / "israel_researcher_overview.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
